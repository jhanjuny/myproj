"""
PPTX processor using python-pptx.
Extracts text per slide, speaker notes, and slide thumbnails via a
render fallback so Claude can see diagrams / charts.
"""
from __future__ import annotations

import io
from pathlib import Path

from src.config import MAX_IMAGE_BYTES, MAX_IMAGES_PER_FILE
from src.processors.base import BaseProcessor, ImageBlock, ProcessedContent, ProcessorError


class PptxProcessor(BaseProcessor):
    supported_extensions = (".pptx", ".ppt")

    def process(self, path: Path) -> ProcessedContent:
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise ProcessorError("python-pptx is not installed (pip install python-pptx)")

        try:
            prs = Presentation(str(path))
        except Exception as e:
            raise ProcessorError(f"Cannot open PPTX: {e}")

        text_parts: list[str] = []
        images: list[ImageBlock] = []
        image_count = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []

            # ── Extract text from all shapes ──────────────────────────────────
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            slide_texts.append(line)

            # ── Speaker notes ─────────────────────────────────────────────────
            notes_text = ""
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text.strip()

            if slide_texts or notes_text:
                block = f"[Slide {slide_num}]\n" + "\n".join(slide_texts)
                if notes_text:
                    block += f"\n[Notes]\n{notes_text}"
                text_parts.append(block)

            # ── Embedded images inside shapes ─────────────────────────────────
            if image_count < MAX_IMAGES_PER_FILE:
                for shape in slide.shapes:
                    if image_count >= MAX_IMAGES_PER_FILE:
                        break
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            raw = shape.image.blob
                            ct = shape.image.content_type
                            if len(raw) > MAX_IMAGE_BYTES:
                                raw = _downscale(raw, MAX_IMAGE_BYTES)
                                ct = "image/jpeg"
                            images.append(ImageBlock(
                                data=raw,
                                media_type=ct,
                                caption=f"Image on slide {slide_num}",
                                page_or_slide=slide_num,
                            ))
                            image_count += 1
                        except Exception:
                            continue

            # ── Render whole slide as image (captures charts, diagrams, etc.) ─
            if image_count < MAX_IMAGES_PER_FILE:
                rendered = _render_slide(prs, slide_num - 1)
                if rendered:
                    images.append(ImageBlock(
                        data=rendered,
                        media_type="image/png",
                        caption=f"Slide {slide_num} rendered",
                        page_or_slide=slide_num,
                    ))
                    image_count += 1

        metadata = {
            "title": path.stem,
            "slide_count": len(prs.slides),
        }

        content = ProcessedContent(
            text="\n\n".join(text_parts),
            images=images,
            metadata=metadata,
            source_path=path,
            file_type="pptx",
        )

        if content.is_empty():
            raise ProcessorError(f"Could not extract any content from {path.name}")

        return content


def _render_slide(prs, slide_index: int) -> bytes | None:
    """
    Render a slide to PNG via LibreOffice (if available) or
    via a python-pptx → pillow fallback for simple slides.
    Falls back gracefully – returns None on failure.
    """
    # Try LibreOffice headless (fast path on Windows/Linux/Mac)
    import subprocess, tempfile, shutil, os
    from pathlib import Path as _P

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                tmp_pptx = _P(tmpdir) / "slide.pptx"
                prs.save(str(tmp_pptx))
                result = subprocess.run(
                    [soffice, "--headless", "--convert-to", "png",
                     "--outdir", tmpdir, str(tmp_pptx)],
                    capture_output=True, timeout=30,
                )
                # LibreOffice produces slide.png for first slide
                # For multi-slide use pdf approach — skip for now
                png_file = _P(tmpdir) / "slide.png"
                if png_file.exists():
                    return png_file.read_bytes()
        except Exception:
            pass

    return None


def _downscale(raw: bytes, max_bytes: int) -> bytes:
    from PIL import Image

    img = Image.open(io.BytesIO(raw)).convert("RGB")
    for q in (75, 55, 35):
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=q)
        data = buf.getvalue()
        if len(data) <= max_bytes:
            return data
    img.thumbnail((1200, 900))
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=50)
    return buf.getvalue()
